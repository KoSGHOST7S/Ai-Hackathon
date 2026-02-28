#!/usr/bin/env python3
"""Generate Assignmint.ai pitch deck as PPTX with faithful extension UI mockups."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from io import BytesIO
from PIL import Image, ImageDraw, ImageFont
import math
import os

# ── Design tokens ──────────────────────────────────────────────
MINT = RGBColor(0x4C, 0xA8, 0x7A)      # #4CA87A - primary
MINT_LIGHT = RGBColor(0xE0, 0xF5, 0xED) # #E0F5ED - accent bg
MINT_BG = RGBColor(0xF4, 0xFA, 0xF7)    # #F4FAF7 - background
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BG = WHITE
FG = RGBColor(0x15, 0x22, 0x28)         # #152228 - foreground
MUTED = RGBColor(0x69, 0x86, 0x7D)      # #69867D - muted text
BORDER = RGBColor(0xB3, 0xD4, 0xC0)     # #B3D4C0 - borders
LIGHT_MUTED = RGBColor(0xE8, 0xF7, 0xF1) # muted bg
ORANGE = RGBColor(0xF4, 0xAE, 0x6C)     # #F4AE6C - secondary
EMERALD = RGBColor(0x10, 0xB9, 0x81)    # emerald for scores
BLUE = RGBColor(0x3B, 0x82, 0xF6)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
ROSE = RGBColor(0xF4, 0x3F, 0x5E)
RED_DESTR = RGBColor(0xE4, 0x43, 0x3A)
CREAM = RGBColor(0xFA, 0xF9, 0xF6)
WARM_BG = RGBColor(0xF3, 0xED, 0xE6)
DIM_TEXT = RGBColor(0xB5, 0xB0, 0xA8)

SLIDE_W = Inches(13.333)  # 16:9 widescreen
SLIDE_H = Inches(7.5)

FONT_DISPLAY = "DM Serif Display"
FONT_BODY = "DM Sans"

# Fallback fonts if DM fonts aren't available
try:
    from pptx.oxml.ns import qn
except ImportError:
    pass


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_name=FONT_BODY,
                font_size=14, color=FG, bold=False, italic=False,
                alignment=PP_ALIGN.LEFT, line_spacing=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = alignment
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    return txBox


def add_rich_textbox(slide, left, top, width, height, runs,
                     alignment=PP_ALIGN.LEFT, line_spacing=None):
    """Add textbox with multiple styled runs. runs = [(text, {font_name, font_size, color, bold, italic}), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    for i, (text, style) in enumerate(runs):
        if i == 0:
            run = p.runs[0] if p.runs else p.add_run()
            run.text = text
        else:
            run = p.add_run()
            run.text = text
        run.font.name = style.get("font_name", FONT_BODY)
        run.font.size = Pt(style.get("font_size", 14))
        run.font.color.rgb = style.get("color", FG)
        run.font.bold = style.get("bold", False)
        run.font.italic = style.get("italic", False)
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width or 1)
    else:
        shape.line.fill.background()
    # Smaller corner radius
    shape.adjustments[0] = 0.08
    return shape


def add_circle(slide, left, top, size, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_line(slide, x1, y1, x2, y2, color=BORDER, width=1):
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)  # MSO_CONNECTOR.STRAIGHT
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    return connector


def create_phone_frame(slide, left, top, width, height):
    """Draw the 390x600 extension popup shell."""
    # Shadow/outer glow (larger rect behind)
    shadow = add_rounded_rect(slide, left - Inches(0.05), top - Inches(0.05),
                              width + Inches(0.1), height + Inches(0.1),
                              RGBColor(0xE0, 0xE0, 0xE0))
    shadow.adjustments[0] = 0.06

    # Main frame
    frame = add_rounded_rect(slide, left, top, width, height, MINT_BG, BORDER, 1.5)
    frame.adjustments[0] = 0.06
    return frame


def draw_ext_header(slide, left, top, width):
    """Draw the Assignmint.ai header bar."""
    # Header background
    header_h = Inches(0.42)
    add_rounded_rect(slide, left, top, width, header_h, WHITE, BORDER, 0.75)

    # Brand name
    add_rich_textbox(slide, left + Inches(0.15), top + Inches(0.08), Inches(2), Inches(0.3), [
        ("Assign", {"font_name": FONT_DISPLAY, "font_size": 13, "color": FG, "bold": True}),
        ("mint.ai", {"font_name": FONT_DISPLAY, "font_size": 13, "color": MINT, "bold": True}),
    ])

    # Avatar circle
    avatar_size = Inches(0.22)
    add_circle(slide, left + width - Inches(0.38), top + Inches(0.1),
               avatar_size, MINT_LIGHT, MINT)
    add_textbox(slide, left + width - Inches(0.38), top + Inches(0.11),
                avatar_size, avatar_size, "MJ",
                font_size=7, color=MINT, bold=True, alignment=PP_ALIGN.CENTER)

    return header_h


def draw_bottom_nav(slide, left, top, width, active="today"):
    """Draw the bottom navigation bar."""
    nav_h = Inches(0.5)
    add_rounded_rect(slide, left, top, width, nav_h, WHITE, BORDER, 0.75)

    tabs = [("Today", "today"), ("Plan", "plan"), ("Me", "me")]
    tab_w = width / 3

    for i, (label, tid) in enumerate(tabs):
        x = left + Emu(int(tab_w * i))
        is_active = tid == active
        icon_text = {"today": "🌿", "plan": "📅", "me": "👤"}[tid]

        if is_active:
            pill = add_rounded_rect(slide, x + Inches(0.08), top + Inches(0.06),
                                    Emu(int(tab_w)) - Inches(0.16), Inches(0.38),
                                    MINT_LIGHT)
            pill.adjustments[0] = 0.3

        add_textbox(slide, x, top + Inches(0.04), Emu(int(tab_w)), Inches(0.22),
                    icon_text, font_size=11,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, top + Inches(0.26), Emu(int(tab_w)), Inches(0.2),
                    label, font_size=7.5,
                    color=MINT if is_active else MUTED,
                    bold=is_active,
                    alignment=PP_ALIGN.CENTER)

    return nav_h


def draw_assignment_card(slide, left, top, width, name, course, points, due_label,
                         due_variant="muted", has_ai=False, highlighted=False):
    """Draw a single assignment card."""
    card_h = Inches(0.6)
    border_c = MINT if highlighted else BORDER
    bg_c = MINT_LIGHT if highlighted else WHITE
    card = add_rounded_rect(slide, left, top, width, card_h, bg_c, border_c, 1)
    card.adjustments[0] = 0.12

    # Title
    add_textbox(slide, left + Inches(0.12), top + Inches(0.08),
                width - Inches(1.4), Inches(0.25),
                name, font_size=9.5, color=FG, bold=True)

    # AI badge
    if has_ai:
        ai_x = left + width - Inches(1.0)
        ai_badge = add_rounded_rect(slide, ai_x, top + Inches(0.08),
                                    Inches(0.38), Inches(0.18), MINT_LIGHT)
        ai_badge.adjustments[0] = 0.3
        add_textbox(slide, ai_x, top + Inches(0.07), Inches(0.38), Inches(0.18),
                    "✨ AI", font_size=7, color=MINT, bold=True, alignment=PP_ALIGN.CENTER)

    # Due badge
    badge_colors = {"muted": (LIGHT_MUTED, MUTED), "secondary": (RGBColor(0xFE, 0xF3, 0xE2), ORANGE),
                    "default": (MINT_LIGHT, MINT)}
    bg, fg = badge_colors.get(due_variant, (LIGHT_MUTED, MUTED))
    due_x = left + width - Inches(0.58)
    due_badge = add_rounded_rect(slide, due_x, top + Inches(0.08),
                                 Inches(0.48), Inches(0.18), bg)
    due_badge.adjustments[0] = 0.4
    add_textbox(slide, due_x, top + Inches(0.07), Inches(0.48), Inches(0.18),
                due_label, font_size=7, color=fg, bold=True, alignment=PP_ALIGN.CENTER)

    # Course + points
    add_textbox(slide, left + Inches(0.12), top + Inches(0.35),
                width - Inches(0.3), Inches(0.2),
                f"⏱ {course} · {points} pts", font_size=8, color=MUTED)

    return card_h


# ═══════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════

def slide_title(prs):
    """Slide 1: Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, CREAM)

    # Subtle gradient blob (mint)
    blob = add_circle(slide, Inches(8), Inches(-1), Inches(6), MINT_LIGHT)
    blob.fill.fore_color.rgb = MINT_LIGHT

    # Logo square
    logo = add_rounded_rect(slide, Inches(5.9), Inches(1.8), Inches(0.7), Inches(0.7), MINT)
    logo.adjustments[0] = 0.18
    add_textbox(slide, Inches(5.9), Inches(1.85), Inches(0.7), Inches(0.7),
                "🌱", font_size=28, alignment=PP_ALIGN.CENTER)

    # Brand
    add_rich_textbox(slide, Inches(2.5), Inches(2.8), Inches(8.5), Inches(0.8), [
        ("Assign", {"font_name": FONT_DISPLAY, "font_size": 56, "color": FG, "bold": True}),
        ("mint.ai", {"font_name": FONT_DISPLAY, "font_size": 56, "color": MINT, "bold": True}),
    ], alignment=PP_ALIGN.CENTER)

    # Divider
    add_line(slide, Inches(5.8), Inches(3.7), Inches(7.5), Inches(3.7), MINT, 1.5)

    # Tagline
    add_textbox(slide, Inches(3), Inches(3.9), Inches(7.3), Inches(0.5),
                "Your assignments, finally organized.",
                font_size=18, color=MUTED, alignment=PP_ALIGN.CENTER)

    # Powered by badge
    badge = add_rounded_rect(slide, Inches(4.7), Inches(4.7), Inches(3.8), Inches(0.4),
                             WHITE, BORDER, 1)
    badge.adjustments[0] = 0.5
    add_textbox(slide, Inches(4.7), Inches(4.73), Inches(3.8), Inches(0.35),
                "⚡  Powered by IBM watsonx Granite",
                font_size=11, color=MUTED, alignment=PP_ALIGN.CENTER)


def slide_problem(prs):
    """Slide 2: The problem — big stat + pain points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CREAM)

    # Section label
    add_textbox(slide, Inches(0.8), Inches(0.6), Inches(3), Inches(0.3),
                "THE PROBLEM", font_size=10, color=MUTED, bold=True)

    # Big stat
    add_textbox(slide, Inches(0.8), Inches(1.2), Inches(5), Inches(2),
                "47%", font_name=FONT_DISPLAY, font_size=120, color=MINT, bold=True)

    add_textbox(slide, Inches(0.8), Inches(3.5), Inches(5.5), Inches(0.8),
                "of students report missing assignments\ndue to poor organization across platforms.",
                font_size=18, color=MUTED, line_spacing=28)

    # Pain point cards on right
    pain_points = [
        "Assignments scattered across tabs",
        "No unified due-date view",
        "Rubrics buried in PDFs",
        "No idea where to start",
    ]
    for i, pain in enumerate(pain_points):
        y = Inches(1.4 + i * 0.75)
        card = add_rounded_rect(slide, Inches(7.2), y, Inches(5), Inches(0.6),
                                RGBColor(0xFE, 0xF0, 0xEB), RGBColor(0xF4, 0xBB, 0xA3), 1)
        card.adjustments[0] = 0.12
        add_textbox(slide, Inches(7.4), y + Inches(0.13), Inches(0.3), Inches(0.3),
                    "✕", font_size=14, color=RGBColor(0xE8, 0x77, 0x5A), bold=True)
        add_textbox(slide, Inches(7.8), y + Inches(0.15), Inches(4.2), Inches(0.3),
                    pain, font_size=12, color=MUTED)


def slide_pain_visual(prs):
    """Slide 3: Canvas wasn't built to help you plan."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CREAM)

    add_textbox(slide, Inches(0.8), Inches(0.6), Inches(3), Inches(0.3),
                "THE REALITY", font_size=10, color=MUTED, bold=True)

    # Headline
    add_rich_textbox(slide, Inches(0.8), Inches(1.3), Inches(6), Inches(1.5), [
        ("Canvas wasn't built\nto help you ", {"font_name": FONT_DISPLAY, "font_size": 42, "color": FG}),
        ("plan.", {"font_name": FONT_DISPLAY, "font_size": 42, "color": MINT, "italic": True}),
    ])

    add_line(slide, Inches(0.8), Inches(3.2), Inches(1.8), Inches(3.2), MINT, 1.5)

    add_textbox(slide, Inches(0.8), Inches(3.5), Inches(5.5), Inches(1.2),
                "Students juggle multiple courses, each with their own assignments, rubrics, and deadlines. Canvas shows you what's due — but never helps you understand what's actually expected.",
                font_size=14, color=MUTED, line_spacing=22)

    # Browser mockup on right
    bx, by = Inches(7.5), Inches(1.0)
    bw, bh = Inches(5), Inches(5.2)

    # Browser frame
    add_rounded_rect(slide, bx, by, bw, bh, WHITE, BORDER, 1)

    # Browser bar
    add_rounded_rect(slide, bx, by, bw, Inches(0.4), RGBColor(0xF5, 0xF3, 0xF0), BORDER, 0.5)
    # Traffic lights
    for j, c in enumerate([RGBColor(0xFF, 0x5F, 0x57), RGBColor(0xFF, 0xBD, 0x2E), RGBColor(0x28, 0xC8, 0x40)]):
        add_circle(slide, bx + Inches(0.15 + j * 0.18), by + Inches(0.12), Inches(0.1), c)

    # URL bar
    add_rounded_rect(slide, bx + Inches(0.8), by + Inches(0.1), Inches(3.5), Inches(0.22),
                     RGBColor(0xEA, 0xE7, 0xE2))

    # Placeholder assignments in browser
    colors = [RGBColor(0xE7, 0x4C, 0x3C), RGBColor(0xF3, 0x9C, 0x12),
              RGBColor(0x34, 0x98, 0xDB), RGBColor(0x95, 0xA5, 0xA6)]
    for i in range(4):
        cy = by + Inches(0.65 + i * 0.7)
        opacity_alpha = 1.0 if i < 3 else 0.5
        card = add_rounded_rect(slide, bx + Inches(0.2), cy, bw - Inches(0.4), Inches(0.55),
                                WHITE, BORDER, 0.5)
        add_circle(slide, bx + Inches(0.35), cy + Inches(0.18), Inches(0.15), colors[i])
        # Placeholder text lines
        add_rounded_rect(slide, bx + Inches(0.65), cy + Inches(0.12),
                         Inches(2.5 if i % 2 == 0 else 2.0), Inches(0.1),
                         RGBColor(0xE8, 0xE3, 0xDC))
        add_rounded_rect(slide, bx + Inches(0.65), cy + Inches(0.3),
                         Inches(1.5 if i % 2 == 0 else 1.8), Inches(0.08),
                         RGBColor(0xE8, 0xE3, 0xDC))


def slide_solution(prs):
    """Slide 4: The solution statement."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CREAM)

    add_textbox(slide, Inches(0.8), Inches(0.6), Inches(3), Inches(0.3),
                "THE SOLUTION", font_size=10, color=MUTED, bold=True)

    add_rich_textbox(slide, Inches(0.8), Inches(2.0), Inches(11), Inches(2), [
        ("Assignmint reads your\nassignments so you\n", {"font_name": FONT_DISPLAY, "font_size": 52, "color": FG}),
        ("don't have to.", {"font_name": FONT_DISPLAY, "font_size": 52, "color": MINT, "italic": True}),
    ])

    add_line(slide, Inches(0.8), Inches(4.8), Inches(2.2), Inches(4.8), MINT, 1.5)

    add_textbox(slide, Inches(0.8), Inches(5.1), Inches(8), Inches(0.6),
                "A Chrome extension that lives inside Canvas — syncing your assignments, analyzing rubrics with AI, and showing you exactly what's expected.",
                font_size=15, color=MUTED, line_spacing=24)


def slide_entry_point(prs):
    """Slide 5: Canvas button — actual Canvas page with button."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CREAM)

    add_textbox(slide, Inches(0.8), Inches(0.6), Inches(3), Inches(0.3),
                "HOW IT WORKS", font_size=10, color=MUTED, bold=True)

    # Step badge
    step_badge = add_rounded_rect(slide, Inches(0.8), Inches(1.2), Inches(1.4), Inches(0.28),
                                  MINT_LIGHT)
    step_badge.adjustments[0] = 0.4
    add_textbox(slide, Inches(0.8), Inches(1.21), Inches(1.4), Inches(0.28),
                "STEP 1", font_size=8, color=MINT, bold=True, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.8), Inches(1.7), Inches(5), Inches(0.9),
                "One click from\nany Canvas page.",
                font_name=FONT_DISPLAY, font_size=36, color=FG)

    add_textbox(slide, Inches(0.8), Inches(3.0), Inches(4.5), Inches(0.8),
                "The \"Open in Assignmint\" button appears on every assignment page. Click it and your assignment is instantly imported and ready to analyze.",
                font_size=13, color=MUTED, line_spacing=21)

    # Mock Assignmint button
    btn = add_rounded_rect(slide, Inches(0.8), Inches(4.2), Inches(2.2), Inches(0.4), MINT)
    btn.adjustments[0] = 0.25
    add_textbox(slide, Inches(0.8), Inches(4.22), Inches(2.2), Inches(0.4),
                "🌿  Open in Assignmint", font_size=11, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)

    # Browser mockup on right showing Canvas
    bx, by = Inches(6.5), Inches(0.8)
    bw, bh = Inches(6.3), Inches(5.8)

    # Browser frame
    add_rounded_rect(slide, bx, by, bw, bh, WHITE, BORDER, 1)

    # Browser bar
    add_rounded_rect(slide, bx, by, bw, Inches(0.38), RGBColor(0xF3, 0xF3, 0xF3), BORDER, 0.5)
    for j, c in enumerate([RGBColor(0xFF, 0x5F, 0x57), RGBColor(0xFF, 0xBD, 0x2E), RGBColor(0x28, 0xC8, 0x40)]):
        add_circle(slide, bx + Inches(0.12 + j * 0.16), by + Inches(0.12), Inches(0.09), c)

    # Tab
    tab_bg = add_rounded_rect(slide, bx + Inches(0.7), by + Inches(0.06),
                               Inches(2.2), Inches(0.26), WHITE)
    add_textbox(slide, bx + Inches(0.8), by + Inches(0.08), Inches(2), Inches(0.22),
                "Canvas LMS — ECON 101", font_size=8, color=FG)

    # URL bar
    add_rounded_rect(slide, bx, by + Inches(0.38), bw, Inches(0.32),
                     RGBColor(0xFA, 0xFA, 0xFA), BORDER, 0.5)
    url_bar = add_rounded_rect(slide, bx + Inches(0.15), by + Inches(0.42),
                                bw - Inches(0.3), Inches(0.22), RGBColor(0xF0, 0xF0, 0xF0))
    add_textbox(slide, bx + Inches(0.25), by + Inches(0.42), bw - Inches(0.5), Inches(0.22),
                "🔒 bw.instructure.com/courses/12345/assignments/67890",
                font_size=8, color=RGBColor(0x99, 0x99, 0x99))

    # Canvas content
    content_y = by + Inches(0.9)

    # Breadcrumb
    add_textbox(slide, bx + Inches(0.25), content_y, Inches(4), Inches(0.2),
                "ECON 101  ›  Assignments  ›  Final Paper",
                font_size=8, color=RGBColor(0x2D, 0x72, 0xD9))

    # Title
    add_textbox(slide, bx + Inches(0.25), content_y + Inches(0.3), Inches(4), Inches(0.4),
                "Final Paper: Supply and Demand Analysis",
                font_size=18, color=RGBColor(0x2D, 0x3B, 0x45))

    # Divider
    add_line(slide, bx + Inches(0.25), content_y + Inches(0.85),
             bx + bw - Inches(0.25), content_y + Inches(0.85),
             RGBColor(0xE0, 0xE0, 0xE0), 0.5)

    # Meta info
    add_textbox(slide, bx + Inches(0.25), content_y + Inches(1.0), Inches(5), Inches(0.2),
                "Due: Mar 15, 2026      Points: 150      Submitting: a file upload",
                font_size=9, color=RGBColor(0x66, 0x66, 0x66))

    # Description
    add_textbox(slide, bx + Inches(0.25), content_y + Inches(1.4), Inches(4.5), Inches(0.6),
                "Write a 2,500-word analysis of supply and demand dynamics in a market of your choice.",
                font_size=10, color=RGBColor(0x55, 0x55, 0x55), line_spacing=17)

    # The Assignmint button on Canvas page
    canvas_btn = add_rounded_rect(slide, bx + bw - Inches(2.3), content_y + Inches(0.15),
                                   Inches(2.0), Inches(0.4), MINT)
    canvas_btn.adjustments[0] = 0.2
    add_textbox(slide, bx + bw - Inches(2.3), content_y + Inches(0.17),
                Inches(2.0), Inches(0.38),
                "🌿  Open in Assignmint",
                font_size=10, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)


def slide_dashboard(prs):
    """Slide 6: Today dashboard — actual extension UI."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CREAM)

    add_textbox(slide, Inches(0.8), Inches(0.6), Inches(3), Inches(0.3),
                "HOW IT WORKS", font_size=10, color=MUTED, bold=True)

    step_badge = add_rounded_rect(slide, Inches(0.8), Inches(1.2), Inches(1.4), Inches(0.28),
                                  MINT_LIGHT)
    step_badge.adjustments[0] = 0.4
    add_textbox(slide, Inches(0.8), Inches(1.21), Inches(1.4), Inches(0.28),
                "STEP 2", font_size=8, color=MINT, bold=True, alignment=PP_ALIGN.CENTER)

    add_rich_textbox(slide, Inches(0.8), Inches(1.7), Inches(5), Inches(1.2), [
        ("Every assignment.\nEvery class. ", {"font_name": FONT_DISPLAY, "font_size": 36, "color": FG}),
        ("One place.", {"font_name": FONT_DISPLAY, "font_size": 36, "color": MINT}),
    ])

    add_textbox(slide, Inches(0.8), Inches(3.3), Inches(4.5), Inches(0.6),
                "All your Canvas assignments sync automatically — organized by due date, grouped by course. No manual entry.",
                font_size=13, color=MUTED, line_spacing=21)

    # Extension mockup on right
    ex, ey = Inches(7.2), Inches(0.6)
    ew, eh = Inches(4.8), Inches(6.2)

    create_phone_frame(slide, ex, ey, ew, eh)
    header_h = draw_ext_header(slide, ex + Inches(0.08), ey + Inches(0.08), ew - Inches(0.16))

    content_y = ey + Inches(0.58)
    content_w = ew - Inches(0.3)
    cx = ex + Inches(0.15)

    # Greeting
    add_textbox(slide, cx, content_y, content_w, Inches(0.25),
                "Good afternoon, Max", font_size=12, color=FG, bold=True)
    add_textbox(slide, cx, content_y + Inches(0.22), content_w, Inches(0.2),
                "2 assignments due today", font_size=8, color=MUTED)

    # Stats row
    stats_y = content_y + Inches(0.5)
    stat_w = Emu(int((content_w - Inches(0.2)) / 3))
    for i, (val, label) in enumerate([(2, "Due today"), (5, "Upcoming"), (12, "All")]):
        sx = cx + Emu(int(i * (stat_w + Inches(0.1))))
        stat_bg = add_rounded_rect(slide, sx, stats_y, stat_w, Inches(0.42), LIGHT_MUTED)
        stat_bg.adjustments[0] = 0.15
        add_textbox(slide, sx, stats_y + Inches(0.04), stat_w, Inches(0.22),
                    str(val), font_size=16, color=MINT, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, sx, stats_y + Inches(0.25), stat_w, Inches(0.15),
                    label, font_size=7, color=MUTED, alignment=PP_ALIGN.CENTER)

    # Assignment cards
    cards_y = stats_y + Inches(0.55)
    cards_data = [
        ("Final Paper: Supply & Demand", "ECON 101", "150", "Due today", "secondary", True, True),
        ("Lab Report #4: Titration", "CHEM 201", "80", "3d left", "default", True, False),
        ("Reading Response Week 8", "ENG 102", "25", "4d left", "muted", False, False),
        ("Problem Set 7: Derivatives", "MATH 301", "60", "5d left", "muted", False, False),
        ("Research Proposal Draft", "PSYC 210", "100", "Due today", "secondary", False, False),
    ]

    for i, (name, course, pts, due, variant, has_ai, highlighted) in enumerate(cards_data):
        cy = cards_y + Emu(int(i * Inches(0.68)))
        if cy + Inches(0.6) > ey + eh - Inches(0.6):
            break
        draw_assignment_card(slide, cx, cy, content_w, name, course, pts, due, variant, has_ai, highlighted)

    # Bottom nav
    draw_bottom_nav(slide, ex + Inches(0.08), ey + eh - Inches(0.55), ew - Inches(0.16), "today")


def slide_ai_analysis(prs):
    """Slide 7: AI analysis pipeline."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CREAM)

    add_textbox(slide, Inches(0.8), Inches(0.6), Inches(3), Inches(0.3),
                "HOW IT WORKS", font_size=10, color=MUTED, bold=True)

    step_badge = add_rounded_rect(slide, Inches(0.8), Inches(1.2), Inches(1.4), Inches(0.28),
                                  MINT_LIGHT)
    step_badge.adjustments[0] = 0.4
    add_textbox(slide, Inches(0.8), Inches(1.21), Inches(1.4), Inches(0.28),
                "STEP 3", font_size=8, color=MINT, bold=True, alignment=PP_ALIGN.CENTER)

    add_rich_textbox(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(1.2), [
        ("Understands rubrics,\nrequirements, and\n", {"font_name": FONT_DISPLAY, "font_size": 36, "color": FG}),
        ("expectations.", {"font_name": FONT_DISPLAY, "font_size": 36, "color": MINT}),
    ])

    add_textbox(slide, Inches(0.8), Inches(3.5), Inches(4.5), Inches(0.8),
                "Powered by IBM watsonx Granite, Assignmint parses your assignment description and rubric — extracting criteria, weights, and actionable steps.",
                font_size=13, color=MUTED, line_spacing=21)

    # Extension mockup showing analysis pipeline
    ex, ey = Inches(7.2), Inches(0.6)
    ew, eh = Inches(4.8), Inches(6.2)

    create_phone_frame(slide, ex, ey, ew, eh)

    # Header with back button
    header_h = Inches(0.42)
    add_rounded_rect(slide, ex, ey, ew, header_h, WHITE, BORDER, 0.75)
    add_textbox(slide, ex + Inches(0.12), ey + Inches(0.08), Inches(0.3), Inches(0.25),
                "←", font_size=12, color=MUTED)
    add_textbox(slide, ex + Inches(0.4), ey + Inches(0.12), Inches(2), Inches(0.2),
                "ECON 101", font_size=9, color=MUTED)

    content_y = ey + Inches(0.6)
    cx = ex + Inches(0.2)
    content_w = ew - Inches(0.4)

    # Assignment title
    add_textbox(slide, cx, content_y, content_w, Inches(0.25),
                "Final Paper: Supply & Demand Analysis",
                font_size=11, color=FG, bold=True)

    # Badges
    badges_y = content_y + Inches(0.3)
    for i, label in enumerate(["150 pts", "Due Mar 15", "online_upload"]):
        bx_pos = cx + Emu(int(i * Inches(0.95)))
        b = add_rounded_rect(slide, bx_pos, badges_y, Inches(0.82), Inches(0.2), LIGHT_MUTED)
        b.adjustments[0] = 0.4
        add_textbox(slide, bx_pos, badges_y + Inches(0.01), Inches(0.82), Inches(0.18),
                    label, font_size=7, color=MUTED, alignment=PP_ALIGN.CENTER)

    # Analysis steps
    steps = [
        ("Extracting explicit requirements…", "done"),
        ("Generating rubric…", "done"),
        ("Validating rubric…", "done"),
        ("Building requirement-based milestones…", "active"),
        ("Validating milestone coverage…", "pending"),
    ]

    steps_y = content_y + Inches(0.8)
    for i, (label, state) in enumerate(steps):
        sy = steps_y + Emu(int(i * Inches(0.45)))
        icon_size = Inches(0.2)

        if state == "done":
            icon = add_circle(slide, cx + Inches(0.1), sy + Inches(0.03), icon_size, MINT)
            add_textbox(slide, cx + Inches(0.1), sy + Inches(0.02), icon_size, icon_size,
                        "✓", font_size=9, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
            text_color = MUTED
        elif state == "active":
            icon = add_circle(slide, cx + Inches(0.1), sy + Inches(0.03), icon_size, MINT_LIGHT, MINT)
            add_textbox(slide, cx + Inches(0.1), sy + Inches(0.02), icon_size, icon_size,
                        "◌", font_size=9, color=MINT, bold=True, alignment=PP_ALIGN.CENTER)
            text_color = FG
        else:
            icon = add_circle(slide, cx + Inches(0.1), sy + Inches(0.03), icon_size, LIGHT_MUTED, BORDER)
            text_color = DIM_TEXT

        add_textbox(slide, cx + Inches(0.42), sy + Inches(0.03), content_w - Inches(0.5), Inches(0.25),
                    label, font_size=10, color=text_color,
                    bold=(state == "active"))

    draw_bottom_nav(slide, ex + Inches(0.08), ey + eh - Inches(0.55), ew - Inches(0.16), "today")


def slide_insights(prs):
    """Slide 8: Results — rubric + milestones."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CREAM)

    add_textbox(slide, Inches(0.8), Inches(0.6), Inches(3), Inches(0.3),
                "HOW IT WORKS", font_size=10, color=MUTED, bold=True)

    step_badge = add_rounded_rect(slide, Inches(0.8), Inches(1.2), Inches(1.4), Inches(0.28),
                                  MINT_LIGHT)
    step_badge.adjustments[0] = 0.4
    add_textbox(slide, Inches(0.8), Inches(1.21), Inches(1.4), Inches(0.28),
                "STEP 4", font_size=8, color=MINT, bold=True, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.8), Inches(1.7), Inches(5), Inches(0.9),
                "Know exactly\nwhat's expected.",
                font_name=FONT_DISPLAY, font_size=36, color=FG)

    add_textbox(slide, Inches(0.8), Inches(3.0), Inches(4.5), Inches(0.8),
                "See your rubric broken down into clear criteria with point weights. Get AI-suggested milestones — a step-by-step path to completing the assignment on time.",
                font_size=13, color=MUTED, line_spacing=21)

    # Extension mockup
    ex, ey = Inches(7.2), Inches(0.6)
    ew, eh = Inches(4.8), Inches(6.2)

    create_phone_frame(slide, ex, ey, ew, eh)

    # Header
    header_h = Inches(0.42)
    add_rounded_rect(slide, ex, ey, ew, header_h, WHITE, BORDER, 0.75)
    add_textbox(slide, ex + Inches(0.12), ey + Inches(0.08), Inches(0.3), Inches(0.25),
                "←", font_size=12, color=MUTED)
    add_textbox(slide, ex + Inches(0.4), ey + Inches(0.12), Inches(2), Inches(0.2),
                "ECON 101", font_size=9, color=MUTED)
    add_textbox(slide, ex + ew - Inches(0.4), ey + Inches(0.12), Inches(0.3), Inches(0.2),
                "✨", font_size=10, alignment=PP_ALIGN.RIGHT)

    cy = ey + Inches(0.55)
    cx = ex + Inches(0.15)
    cw = ew - Inches(0.3)

    # Rubric section
    add_textbox(slide, cx, cy, Inches(1), Inches(0.2),
                "RUBRIC", font_size=7, color=MUTED, bold=True)
    add_textbox(slide, cx + cw - Inches(0.5), cy, Inches(0.5), Inches(0.2),
                "150 pts", font_size=7, color=MUTED, alignment=PP_ALIGN.RIGHT)

    criteria = [
        ("Content & Argumentation", "40"),
        ("Research Quality", "35"),
        ("Structure & Organization", "30"),
        ("Evidence & Analysis", "25"),
        ("Citations & Formatting", "20"),
    ]

    for i, (name, weight) in enumerate(criteria):
        ry = cy + Inches(0.25 + i * 0.38)
        card = add_rounded_rect(slide, cx, ry, cw, Inches(0.32), WHITE, BORDER, 0.5)
        card.adjustments[0] = 0.12
        add_textbox(slide, cx + Inches(0.1), ry + Inches(0.06), cw - Inches(0.9), Inches(0.2),
                    name, font_size=8.5, color=FG, bold=True)
        # Weight badge
        wb = add_rounded_rect(slide, cx + cw - Inches(0.55), ry + Inches(0.06),
                               Inches(0.35), Inches(0.18), LIGHT_MUTED)
        wb.adjustments[0] = 0.4
        add_textbox(slide, cx + cw - Inches(0.55), ry + Inches(0.05), Inches(0.35), Inches(0.18),
                    weight, font_size=7, color=MUTED, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, cx + cw - Inches(0.18), ry + Inches(0.06), Inches(0.15), Inches(0.18),
                    "›", font_size=10, color=MUTED)

    # Milestones section
    my = cy + Inches(2.2)
    add_textbox(slide, cx, my, Inches(1.5), Inches(0.2),
                "MILESTONES", font_size=7, color=MUTED, bold=True)
    add_textbox(slide, cx + cw - Inches(1), my, Inches(1), Inches(0.2),
                "1/5 complete", font_size=7, color=MUTED, alignment=PP_ALIGN.RIGHT)

    milestones = [
        ("Choose market & gather data", "~3h", True),
        ("Draft elasticity analysis", "~4h", False),
        ("Write policy implications", "~3h", False),
    ]

    for i, (title, hours, checked) in enumerate(milestones):
        miy = my + Inches(0.25 + i * 0.42)
        card = add_rounded_rect(slide, cx, miy, cw, Inches(0.36), WHITE, BORDER, 0.5)
        card.adjustments[0] = 0.12

        # Checkbox
        if checked:
            chk = add_circle(slide, cx + Inches(0.1), miy + Inches(0.08), Inches(0.16), MINT)
            add_textbox(slide, cx + Inches(0.1), miy + Inches(0.06), Inches(0.16), Inches(0.16),
                        "✓", font_size=7, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        else:
            chk = add_rounded_rect(slide, cx + Inches(0.1), miy + Inches(0.08),
                                    Inches(0.16), Inches(0.16), LIGHT_MUTED, BORDER, 0.5)
            add_textbox(slide, cx + Inches(0.1), miy + Inches(0.07), Inches(0.16), Inches(0.16),
                        str(i + 1), font_size=6, color=MUTED, bold=True, alignment=PP_ALIGN.CENTER)

        add_textbox(slide, cx + Inches(0.35), miy + Inches(0.06), cw - Inches(0.8), Inches(0.15),
                    title, font_size=8.5, color=MUTED if checked else FG, bold=not checked)
        add_textbox(slide, cx + Inches(0.35), miy + Inches(0.2), cw - Inches(0.8), Inches(0.12),
                    f"⏱ {hours}", font_size=7, color=MUTED)

    draw_bottom_nav(slide, ex + Inches(0.08), ey + eh - Inches(0.55), ew - Inches(0.16), "today")


def slide_chat(prs):
    """Slide 9: AI Chat."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CREAM)

    add_textbox(slide, Inches(0.8), Inches(0.6), Inches(3), Inches(0.3),
                "HOW IT WORKS", font_size=10, color=MUTED, bold=True)

    step_badge = add_rounded_rect(slide, Inches(0.8), Inches(1.2), Inches(1.4), Inches(0.28),
                                  MINT_LIGHT)
    step_badge.adjustments[0] = 0.4
    add_textbox(slide, Inches(0.8), Inches(1.21), Inches(1.4), Inches(0.28),
                "STEP 5", font_size=8, color=MINT, bold=True, alignment=PP_ALIGN.CENTER)

    add_rich_textbox(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(1.5), [
        ("Ask anything.\nGet answers grounded\nin ", {"font_name": FONT_DISPLAY, "font_size": 36, "color": FG}),
        ("the actual rubric.", {"font_name": FONT_DISPLAY, "font_size": 36, "color": MINT}),
    ])

    add_textbox(slide, Inches(0.8), Inches(3.6), Inches(4.5), Inches(0.8),
                "Not generic ChatGPT. Every response is grounded in your specific assignment, rubric, and course context.",
                font_size=13, color=MUTED, line_spacing=21)

    # Extension mockup
    ex, ey = Inches(7.2), Inches(0.6)
    ew, eh = Inches(4.8), Inches(6.2)

    create_phone_frame(slide, ex, ey, ew, eh)

    # Header
    header_h = Inches(0.42)
    add_rounded_rect(slide, ex, ey, ew, header_h, WHITE, BORDER, 0.75)
    add_textbox(slide, ex + Inches(0.12), ey + Inches(0.08), Inches(0.3), Inches(0.25),
                "←", font_size=12, color=MUTED)
    add_textbox(slide, ex + Inches(0.4), ey + Inches(0.12), Inches(1), Inches(0.2),
                "Ask AI", font_size=10, color=FG, bold=True)

    cx = ex + Inches(0.15)
    cw = ew - Inches(0.3)
    cy = ey + Inches(0.6)

    # Bot intro message
    bot_icon_size = Inches(0.2)
    add_circle(slide, cx, cy, bot_icon_size, MINT_LIGHT)
    add_textbox(slide, cx, cy - Inches(0.01), bot_icon_size, bot_icon_size,
                "🤖", font_size=8, alignment=PP_ALIGN.CENTER)

    intro_bubble = add_rounded_rect(slide, cx + Inches(0.28), cy, cw - Inches(0.4), Inches(0.52), LIGHT_MUTED)
    intro_bubble.adjustments[0] = 0.12
    add_textbox(slide, cx + Inches(0.35), cy + Inches(0.06), cw - Inches(0.55), Inches(0.45),
                "I can answer questions about \"Final Paper: Supply & Demand Analysis\" — its rubric, milestones, and requirements. What would you like to know?",
                font_size=8, color=FG, line_spacing=13)

    # User message
    user_y = cy + Inches(0.65)
    user_bubble = add_rounded_rect(slide, cx + cw - Inches(2.8), user_y,
                                    Inches(2.5), Inches(0.32), MINT)
    user_bubble.adjustments[0] = 0.15
    add_textbox(slide, cx + cw - Inches(2.7), user_y + Inches(0.06),
                Inches(2.3), Inches(0.22),
                "What makes a strong thesis for this paper?",
                font_size=8, color=WHITE)

    user_icon = add_circle(slide, cx + cw - Inches(0.22), user_y + Inches(0.05),
                           bot_icon_size, LIGHT_MUTED)
    add_textbox(slide, cx + cw - Inches(0.22), user_y + Inches(0.04),
                bot_icon_size, bot_icon_size,
                "👤", font_size=8, alignment=PP_ALIGN.CENTER)

    # Bot response
    resp_y = user_y + Inches(0.45)
    add_circle(slide, cx, resp_y, bot_icon_size, MINT_LIGHT)
    add_textbox(slide, cx, resp_y - Inches(0.01), bot_icon_size, bot_icon_size,
                "🤖", font_size=8, alignment=PP_ALIGN.CENTER)

    resp_bubble = add_rounded_rect(slide, cx + Inches(0.28), resp_y,
                                    cw - Inches(0.4), Inches(1.8), LIGHT_MUTED)
    resp_bubble.adjustments[0] = 0.08
    add_textbox(slide, cx + Inches(0.35), resp_y + Inches(0.06),
                cw - Inches(0.55), Inches(1.7),
                "Based on the rubric's \"Content & Argumentation\" criteria (40 pts), a strong thesis should:\n\n"
                "1. State a clear market choice\n"
                "2. Include a debatable claim about supply/demand dynamics\n"
                "3. Preview your policy implications\n\n"
                "Example: \"The U.S. housing market demonstrates how inelastic supply combined with policy-driven demand creates persistent affordability challenges...\"",
                font_size=8, color=FG, line_spacing=13)

    # Input bar at bottom
    input_y = ey + eh - Inches(1.0)
    add_line(slide, cx, input_y, cx + cw, input_y, BORDER, 0.5)
    input_bar = add_rounded_rect(slide, cx, input_y + Inches(0.08),
                                  cw - Inches(0.4), Inches(0.3), LIGHT_MUTED)
    input_bar.adjustments[0] = 0.2
    add_textbox(slide, cx + Inches(0.1), input_y + Inches(0.1),
                cw - Inches(0.6), Inches(0.25),
                "Ask about this assignment…", font_size=8, color=DIM_TEXT)

    send_btn = add_rounded_rect(slide, cx + cw - Inches(0.35), input_y + Inches(0.08),
                                 Inches(0.3), Inches(0.3), MINT)
    send_btn.adjustments[0] = 0.15
    add_textbox(slide, cx + cw - Inches(0.35), input_y + Inches(0.08),
                Inches(0.3), Inches(0.3),
                "▸", font_size=12, color=WHITE, alignment=PP_ALIGN.CENTER)


def slide_score(prs):
    """Slide 10: Pre-submission score / review."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CREAM)

    add_textbox(slide, Inches(0.8), Inches(0.6), Inches(3), Inches(0.3),
                "HOW IT WORKS", font_size=10, color=MUTED, bold=True)

    step_badge = add_rounded_rect(slide, Inches(0.8), Inches(1.2), Inches(1.4), Inches(0.28),
                                  MINT_LIGHT)
    step_badge.adjustments[0] = 0.4
    add_textbox(slide, Inches(0.8), Inches(1.21), Inches(1.4), Inches(0.28),
                "STEP 6", font_size=8, color=MINT, bold=True, alignment=PP_ALIGN.CENTER)

    add_rich_textbox(slide, Inches(0.8), Inches(1.7), Inches(5), Inches(0.9), [
        ("Know your grade\n", {"font_name": FONT_DISPLAY, "font_size": 36, "color": FG}),
        ("before you submit.", {"font_name": FONT_DISPLAY, "font_size": 36, "color": MINT}),
    ])

    add_textbox(slide, Inches(0.8), Inches(3.0), Inches(4.5), Inches(0.8),
                "Upload your draft and get an AI-predicted score against the rubric. See per-category breakdowns and the single most impactful improvement.",
                font_size=13, color=MUTED, line_spacing=21)

    # Extension mockup
    ex, ey = Inches(7.2), Inches(0.6)
    ew, eh = Inches(4.8), Inches(6.2)

    create_phone_frame(slide, ex, ey, ew, eh)

    # Header
    header_h = Inches(0.42)
    add_rounded_rect(slide, ex, ey, ew, header_h, WHITE, BORDER, 0.75)
    add_textbox(slide, ex + Inches(0.12), ey + Inches(0.08), Inches(0.3), Inches(0.25),
                "←", font_size=12, color=MUTED)
    add_rich_textbox(slide, ex + Inches(0.4), ey + Inches(0.1), Inches(2), Inches(0.22), [
        ("Assign", {"font_name": FONT_DISPLAY, "font_size": 10, "color": FG, "bold": True}),
        ("mint.ai", {"font_name": FONT_DISPLAY, "font_size": 10, "color": MINT, "bold": True}),
        (" · ", {"font_size": 10, "color": BORDER}),
        ("AI Review", {"font_size": 10, "color": FG, "bold": True}),
    ])

    cx = ex + Inches(0.15)
    cw = ew - Inches(0.3)
    cy = ey + Inches(0.6)

    # Score display
    score_card = add_rounded_rect(slide, cx, cy, cw, Inches(1.2), WHITE, BORDER, 0.75)
    score_card.adjustments[0] = 0.08

    add_textbox(slide, cx, cy + Inches(0.08), cw, Inches(0.15),
                "PRE-SUBMISSION SCORE", font_size=7, color=MUTED, bold=True,
                alignment=PP_ALIGN.CENTER)

    # Big score number
    add_rich_textbox(slide, cx, cy + Inches(0.28), cw, Inches(0.55), [
        ("128", {"font_name": FONT_DISPLAY, "font_size": 40, "color": MINT, "bold": True}),
        ("/150", {"font_size": 16, "color": MUTED}),
    ], alignment=PP_ALIGN.CENTER)

    # Percentage
    pct_y = cy + Inches(0.85)
    dot = add_circle(slide, cx + Emu(int(cw / 2)) - Inches(0.45), pct_y, Inches(0.1), EMERALD)
    add_textbox(slide, cx + Emu(int(cw / 2)) - Inches(0.32), pct_y - Inches(0.02),
                Inches(1.5), Inches(0.18),
                "85% — Strong pass", font_size=8.5, color=MUTED)

    # Breakdown section
    by = cy + Inches(1.35)
    add_textbox(slide, cx, by, Inches(1.5), Inches(0.2),
                "BREAKDOWN", font_size=7, color=MUTED, bold=True)

    scores = [
        ("Content & Argumentation", 36, 40, EMERALD),
        ("Research Quality", 31, 35, MINT),
        ("Structure & Organization", 26, 30, MINT),
        ("Evidence & Analysis", 21, 25, MINT),
        ("Citations & Formatting", 14, 20, AMBER),
    ]

    for i, (name, score, max_score, bar_color) in enumerate(scores):
        ry = by + Inches(0.25 + i * 0.42)

        # Name and score
        add_textbox(slide, cx + Inches(0.05), ry, cw - Inches(0.6), Inches(0.15),
                    name, font_size=8, color=FG, bold=True)
        add_textbox(slide, cx + cw - Inches(0.55), ry, Inches(0.5), Inches(0.15),
                    f"{score}/{max_score}", font_size=8, color=MINT, bold=True,
                    alignment=PP_ALIGN.RIGHT)

        # Bar background
        bar_y = ry + Inches(0.18)
        bar_bg = add_rounded_rect(slide, cx + Inches(0.05), bar_y, cw - Inches(0.1), Inches(0.06), LIGHT_MUTED)
        bar_bg.adjustments[0] = 0.5

        # Bar fill
        fill_w = Emu(int((cw - Inches(0.1)) * (score / max_score)))
        if fill_w > 0:
            bar_fill = add_rounded_rect(slide, cx + Inches(0.05), bar_y, fill_w, Inches(0.06), bar_color)
            bar_fill.adjustments[0] = 0.5

    # Improvement tip
    tip_y = by + Inches(2.35)
    tip_card = add_rounded_rect(slide, cx, tip_y, cw, Inches(0.65),
                                 MINT_LIGHT, RGBColor(0x8B, 0xCC, 0xA6), 0.75)
    tip_card.adjustments[0] = 0.1
    add_textbox(slide, cx + Inches(0.1), tip_y + Inches(0.06), cw - Inches(0.2), Inches(0.12),
                "TOP IMPROVEMENT", font_size=7, color=MINT, bold=True)
    add_textbox(slide, cx + Inches(0.1), tip_y + Inches(0.22), cw - Inches(0.2), Inches(0.4),
                "Citations could be more consistent. Switch to one citation style and ensure all in-text citations match your bibliography.",
                font_size=8, color=FG, line_spacing=13)

    draw_bottom_nav(slide, ex + Inches(0.08), ey + eh - Inches(0.55), ew - Inches(0.16), "today")


def slide_tech_stack(prs):
    """Slide 11: Tech stack."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CREAM)

    add_textbox(slide, Inches(0.8), Inches(0.6), Inches(3), Inches(0.3),
                "UNDER THE HOOD", font_size=10, color=MUTED, bold=True)

    add_rich_textbox(slide, Inches(0.8), Inches(1.5), Inches(10), Inches(1), [
        ("Built to be fast, private,\nand ", {"font_name": FONT_DISPLAY, "font_size": 42, "color": FG}),
        ("reliable.", {"font_name": FONT_DISPLAY, "font_size": 42, "color": MINT}),
    ])

    add_textbox(slide, Inches(0.8), Inches(3.0), Inches(6), Inches(0.4),
                "Lightweight Chrome extension with enterprise-grade AI.",
                font_size=16, color=MUTED)

    # Tech cards
    tech = [
        ("🧩", "Chrome Extension", "Manifest V3, runs locally"),
        ("⚛️", "React + TypeScript", "Type-safe, component-driven"),
        ("🤖", "IBM watsonx Granite", "Rubric analysis & scoring"),
        ("🔗", "Canvas LMS API", "Auto-sync assignments & rubrics"),
    ]

    for i, (icon, name, desc) in enumerate(tech):
        col = i % 2
        row = i // 2
        x = Inches(0.8 + col * 5.8)
        y = Inches(3.8 + row * 1.2)

        card = add_rounded_rect(slide, x, y, Inches(5.2), Inches(0.9), WHITE, BORDER, 1)
        card.adjustments[0] = 0.1

        # Icon circle
        icon_bg_colors = [MINT_LIGHT, RGBColor(0xE3, 0xF2, 0xFD),
                          RGBColor(0xFF, 0xF3, 0xE0), RGBColor(0xFE, 0xF0, 0xEB)]
        icon_circle = add_rounded_rect(slide, x + Inches(0.2), y + Inches(0.2),
                                        Inches(0.5), Inches(0.5), icon_bg_colors[i])
        icon_circle.adjustments[0] = 0.2
        add_textbox(slide, x + Inches(0.2), y + Inches(0.22), Inches(0.5), Inches(0.5),
                    icon, font_size=18, alignment=PP_ALIGN.CENTER)

        add_textbox(slide, x + Inches(0.9), y + Inches(0.18), Inches(3.5), Inches(0.25),
                    name, font_size=13, color=FG, bold=True)
        add_textbox(slide, x + Inches(0.9), y + Inches(0.45), Inches(3.5), Inches(0.25),
                    desc, font_size=10, color=MUTED)


def slide_cta(prs):
    """Slide 12: CTA."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CREAM)

    # Blob
    blob = add_circle(slide, Inches(4), Inches(0.5), Inches(8), MINT_LIGHT)

    # Logo
    logo = add_rounded_rect(slide, Inches(6.0), Inches(1.5), Inches(0.8), Inches(0.8), MINT)
    logo.adjustments[0] = 0.18
    add_textbox(slide, Inches(6.0), Inches(1.55), Inches(0.8), Inches(0.8),
                "🌱", font_size=32, alignment=PP_ALIGN.CENTER)

    # Brand
    add_rich_textbox(slide, Inches(2.5), Inches(2.6), Inches(8.5), Inches(0.8), [
        ("Assign", {"font_name": FONT_DISPLAY, "font_size": 52, "color": FG, "bold": True}),
        ("mint.ai", {"font_name": FONT_DISPLAY, "font_size": 52, "color": MINT, "bold": True}),
    ], alignment=PP_ALIGN.CENTER)

    add_line(slide, Inches(5.8), Inches(3.5), Inches(7.5), Inches(3.5), MINT, 1.5)

    add_textbox(slide, Inches(3), Inches(3.7), Inches(7.3), Inches(0.5),
                "Stop guessing. Start knowing.",
                font_size=18, color=MUTED, alignment=PP_ALIGN.CENTER)

    # CTA button
    cta_btn = add_rounded_rect(slide, Inches(5.0), Inches(4.5), Inches(3.3), Inches(0.55), MINT)
    cta_btn.adjustments[0] = 0.2
    add_textbox(slide, Inches(5.0), Inches(4.52), Inches(3.3), Inches(0.52),
                "Try Assignmint.ai  →", font_size=15, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)

    # Footer links
    add_textbox(slide, Inches(2), Inches(5.5), Inches(9.3), Inches(0.3),
                "Chrome Web Store   ·   assignmint.ai   ·   Built with IBM watsonx",
                font_size=11, color=DIM_TEXT, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_problem(prs)
    slide_pain_visual(prs)
    slide_solution(prs)
    slide_entry_point(prs)
    slide_dashboard(prs)
    slide_ai_analysis(prs)
    slide_insights(prs)
    slide_chat(prs)
    slide_score(prs)
    slide_tech_stack(prs)
    slide_cta(prs)

    out_path = os.path.join(os.path.dirname(__file__), "assignmint-pitch-deck.pptx")
    prs.save(out_path)
    print(f"Saved: {out_path}")


if __name__ == "__main__":
    main()
